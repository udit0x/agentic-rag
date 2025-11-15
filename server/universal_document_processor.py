"""Universal document processing pipeline supporting multiple file types."""
import logging
import base64
import io
import json
import tempfile
import os
from pathlib import Path
from typing import Tuple, Dict, Any, Optional
import mimetypes

# Configuration management
from server.config_manager import config_manager

# PDF processing
import fitz  # pymupdf

logger = logging.getLogger(__name__)

# Document loaders
try:
    from langchain_community.document_loaders import (
        UnstructuredWordDocumentLoader,
        UnstructuredPowerPointLoader,
        UnstructuredMarkdownLoader
    )
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    logger.warning("Unstructured not available, falling back to alternative methods")
    UNSTRUCTURED_AVAILABLE = False

# CSV/Excel processing
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    logger.warning("Pandas not available, CSV support limited")
    PANDAS_AVAILABLE = False

# DOCX fallback
try:
    from docx import Document as DocxDocument
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    logger.warning("python-docx not available")
    PYTHON_DOCX_AVAILABLE = False

# PPTX fallback
try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    logger.warning("python-pptx not available")
    PYTHON_PPTX_AVAILABLE = False

# OCR support
try:
    from .ocr_processor import ocr_processor
except ImportError:
    from ocr_processor import ocr_processor


def detect_file_type(content_type: str, filename: str) -> str:
    """
    Detect file type from MIME type and filename extension.
    
    Args:
        content_type: MIME type of the file
        filename: Original filename
        
    Returns:
        Standardized file type string
    """
    ext = Path(filename).suffix.lower()
    mime = content_type.lower()
    
    # PDF
    if "pdf" in mime or ext == ".pdf":
        return "pdf"
    
    # Microsoft Word
    elif ("word" in mime or "msword" in mime or 
          ext in [".doc", ".docx"] or
          "officedocument.wordprocessingml" in mime):
        return "docx"
    
    # Microsoft PowerPoint
    elif ("presentation" in mime or "powerpoint" in mime or
          ext in [".ppt", ".pptx"] or
          "officedocument.presentationml" in mime):
        return "pptx"
    
    # CSV
    elif "csv" in mime or ext == ".csv":
        return "csv"
    
    # Excel
    elif ("excel" in mime or "spreadsheet" in mime or
          ext in [".xls", ".xlsx"] or
          "officedocument.spreadsheetml" in mime):
        return "xlsx"
    
    # JSON
    elif "json" in mime or ext == ".json":
        return "json"
    
    # Markdown
    elif "markdown" in mime or ext in [".md", ".markdown"]:
        return "markdown"
    
    # HTML
    elif "html" in mime or ext in [".html", ".htm"]:
        return "html"
    
    # Plain text
    elif "text" in mime or ext == ".txt":
        return "txt"
    
    else:
        # Try to guess from filename if MIME type is generic
        guessed_type, _ = mimetypes.guess_type(filename)
        if guessed_type and guessed_type.lower() != content_type.lower():
            # Avoid infinite recursion by checking if the guessed type is different
            return detect_file_type(guessed_type, filename)
        
        raise ValueError(f"Unsupported file type: {mime} ({ext})")


class DocumentExtractor:
    """Universal document text extraction class."""
    
    @staticmethod
    async def extract_pdf_text(file_bytes: bytes, filename: str, force_ocr: bool = False) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF with OCR fallback."""
        logger.info("Processing PDF: %s", filename)
        
        metrics = {
            "extraction_mode": "Direct",
            "ocr_triggered": False,
            "page_count": 0,
            "total_chars": 0,
            "errors": []
        }
        
        try:
            # Extract text using PyMuPDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text_content = ""
            metrics["page_count"] = len(doc)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                text_content += page_text + "\n"
            
            doc.close()
            text_content = text_content.strip()
            metrics["total_chars"] = len(text_content)
            
            # Check if OCR is needed
            ocr_needed = force_ocr or not ocr_processor.is_text_content_sufficient(text_content)
            
            if ocr_needed:
                logger.info("PDF requires OCR processing: %s", filename)
                metrics["ocr_triggered"] = True
                metrics["extraction_mode"] = "OCR"
                
                try:
                    ocr_text, ocr_metrics = await ocr_processor.process_scanned_document(
                        file_bytes, filename
                    )
                    
                    if ocr_text and len(ocr_text.strip()) > 0:
                        text_content = ocr_text
                        metrics.update(ocr_metrics.finalize())
                        logger.info("OCR successful: %d chars extracted", len(text_content))
                    else:
                        metrics["errors"].append("OCR extraction yielded no content")
                        
                except Exception as ocr_error:
                    logger.error("OCR failed for %s: %s", filename, str(ocr_error))
                    metrics["errors"].append(f"OCR failed: {str(ocr_error)}")
            
            return text_content, metrics
            
        except Exception as e:
            error_msg = f"PDF extraction failed: {str(e)}"
            logger.error("PDF extraction failed for %s: %s", filename, str(e))
            metrics["errors"].append(error_msg)
            return "", metrics
    
    @staticmethod
    def extract_docx_text(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX files."""
        logger.info("Processing DOCX: %s", filename)
        
        metrics = {
            "extraction_mode": "Direct",
            "total_chars": 0,
            "errors": []
        }
        
        # Try Unstructured first
        if UNSTRUCTURED_AVAILABLE:
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                
                loader = UnstructuredWordDocumentLoader(tmp_path)
                docs = loader.load()
                text_content = "\n".join([d.page_content for d in docs])
                
                # Cleanup
                os.unlink(tmp_path)
                
                metrics["total_chars"] = len(text_content)
                logger.info("DOCX extracted via Unstructured: %d chars", len(text_content))
                return text_content, metrics
                
            except Exception as e:
                logger.warning("Unstructured DOCX extraction failed for %s: %s", filename, str(e))
                metrics["errors"].append(f"Unstructured failed: {str(e)}")
        
        # Fallback to python-docx
        if PYTHON_DOCX_AVAILABLE:
            try:
                doc = DocxDocument(io.BytesIO(file_bytes))
                text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                
                metrics["total_chars"] = len(text_content)
                logger.info("DOCX extracted via python-docx: %d chars", len(text_content))
                return text_content, metrics
                
            except Exception as e:
                logger.warning("python-docx extraction failed for %s: %s", filename, str(e))
                metrics["errors"].append(f"python-docx failed: {str(e)}")
        
        # If all methods fail
        error_msg = "No DOCX extraction method available"
        logger.error("DOCX extraction failed for %s: %s", filename, error_msg)
        metrics["errors"].append(error_msg)
        return "", metrics
    
    @staticmethod
    def extract_pptx_text(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PPTX files."""
        logger.info("Processing PPTX: %s", filename)
        
        metrics = {
            "extraction_mode": "Direct",
            "total_chars": 0,
            "errors": []
        }
        
        # Try Unstructured first
        if UNSTRUCTURED_AVAILABLE:
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                
                loader = UnstructuredPowerPointLoader(tmp_path)
                docs = loader.load()
                text_content = "\n".join([d.page_content for d in docs])
                
                # Cleanup
                os.unlink(tmp_path)
                
                metrics["total_chars"] = len(text_content)
                logger.info("PPTX extracted via Unstructured: %d chars", len(text_content))
                return text_content, metrics
                
            except Exception as e:
                logger.warning("Unstructured PPTX extraction failed for %s: %s", filename, str(e))
                metrics["errors"].append(f"Unstructured failed: {str(e)}")
        
        # Fallback to python-pptx
        if PYTHON_PPTX_AVAILABLE:
            try:
                prs = Presentation(io.BytesIO(file_bytes))
                text_content = ""
                
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = f"=== Slide {slide_num + 1} ===\n"
                    
                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text += shape.text + "\n"
                    
                    # Extract notes if available
                    if slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text.strip():
                            slide_text += f"Notes: {notes_text}\n"
                    
                    text_content += slide_text + "\n"
                
                metrics["total_chars"] = len(text_content)
                logger.info("PPTX extracted via python-pptx: %d chars", len(text_content))
                return text_content, metrics
                
            except Exception as e:
                logger.warning("python-pptx extraction failed for %s: %s", filename, str(e))
                metrics["errors"].append(f"python-pptx failed: {str(e)}")
        
        # If all methods fail
        error_msg = "No PPTX extraction method available"
        logger.error("PPTX extraction failed for %s: %s", filename, error_msg)
        metrics["errors"].append(error_msg)
        return "", metrics
    
    @staticmethod
    def extract_csv_text(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from CSV files with intelligent sampling."""
        logger.info("Processing CSV: %s", filename)
        
        metrics = {
            "extraction_mode": "Smart Sampling",
            "total_chars": 0,
            "rows_sampled": 0,
            "optimization_applied": False,
            "errors": []
        }
        
        if not PANDAS_AVAILABLE:
            # Fallback to basic CSV reading with sampling
            try:
                import csv
                text_lines = []
                csv_text = file_bytes.decode("utf-8", errors="ignore")
                reader = csv.reader(io.StringIO(csv_text))
                
                lines = list(reader)
                total_rows = len(lines)
                
                if total_rows <= 100:
                    # Small CSV: include everything
                    text_lines = [" | ".join(row) for row in lines]
                    metrics["rows_sampled"] = total_rows
                else:
                    # Large CSV: sample intelligently
                    metrics["optimization_applied"] = True
                    text_lines.append(f"CSV Summary: {total_rows} total rows, showing sample")
                    text_lines.append("")
                    
                    # Header + first 20 rows
                    for i, row in enumerate(lines[:21]):
                        text_lines.append(" | ".join(row))
                    
                    if total_rows > 50:
                        text_lines.append("\n--- Last 10 rows ---")
                        for row in lines[-10:]:
                            text_lines.append(" | ".join(row))
                    
                    metrics["rows_sampled"] = min(31, total_rows)
                
                text_content = "\n".join(text_lines)
                metrics["total_chars"] = len(text_content)
                logger.info("CSV extracted via basic reader: %d total rows, %d sampled, %d chars", 
                           total_rows, metrics['rows_sampled'], len(text_content))
                return text_content, metrics
                
            except Exception as e:
                error_msg = f"Basic CSV extraction failed: {str(e)}"
                logger.error("Basic CSV extraction failed for %s: %s", filename, str(e))
                metrics["errors"].append(error_msg)
                return "", metrics
        
        try:
            df = pd.read_csv(io.BytesIO(file_bytes))
            original_rows = len(df)
            
            # Smart processing based on CSV size
            if original_rows <= 100:
                # Small CSV: include everything
                text_content = f"CSV File: {filename} ({original_rows} rows)\n\n"
                text_content += df.to_string(max_cols=20)
                metrics["rows_sampled"] = original_rows
                
            elif original_rows <= 500:
                # Medium CSV: include with some limits
                text_content = f"CSV File: {filename} ({original_rows} rows)\n\n"
                text_content += df.to_string(max_rows=500, max_cols=15)
                metrics["rows_sampled"] = original_rows
                
            else:
                # Large CSV: intelligent sampling
                metrics["optimization_applied"] = True
                
                text_content = f"CSV File: {filename} ({original_rows:,} rows - sampled)\n"
                text_content += f"Columns: {', '.join(df.columns.astype(str))}\n\n"
                
                # Include first 30 rows (headers + initial data)
                text_content += "--- First 30 rows ---\n"
                text_content += df.head(30).to_string(max_cols=15) + "\n\n"
                
                # Include last 15 rows
                text_content += "--- Last 15 rows ---\n"
                text_content += df.tail(15).to_string(max_cols=15) + "\n\n"
                
                # Sample middle if very large
                if original_rows > 1000:
                    middle_start = original_rows // 2
                    middle_sample = df.iloc[middle_start:middle_start + 15]
                    text_content += f"--- Sample from middle (rows {middle_start}-{middle_start + 15}) ---\n"
                    text_content += middle_sample.to_string(max_cols=15) + "\n\n"
                
                # Add summary statistics for numerical columns
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0 and len(numeric_cols) <= 10:
                    text_content += "--- Summary Statistics ---\n"
                    text_content += df[numeric_cols].describe().round(2).to_string() + "\n\n"
                
                metrics["rows_sampled"] = 60  # 30 + 15 + 15 sampled rows
                metrics["errors"].append(f"Large CSV sampled from {original_rows:,} to {metrics['rows_sampled']} rows")
            
            # Final size check
            if len(text_content) > 50000:  # 50KB limit for CSV
                text_content = text_content[:50000] + "\n\n... [Content truncated for processing efficiency]"
                metrics["errors"].append("Content truncated to 50KB")
            
            metrics["total_chars"] = len(text_content)
            logger.info("CSV processed: %d total rows, %d sampled, %d chars", 
                       original_rows, metrics['rows_sampled'], len(text_content))
            return text_content, metrics
            
        except Exception as e:
            error_msg = f"Pandas CSV extraction failed: {str(e)}"
            logger.error("Pandas CSV extraction failed for %s: %s", filename, str(e))
            metrics["errors"].append(error_msg)
            return "", metrics
    
    @staticmethod
    def extract_xlsx_text(file_bytes: bytes, filename: str, use_smart_sampling: bool = True) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Excel files with intelligent sampling and summarization."""
        logger.info("Processing XLSX: %s (Smart Sampling: %s)", filename, use_smart_sampling)
        
        metrics = {
            "extraction_mode": "Smart Sampling",
            "total_chars": 0,
            "sheets_processed": 0,
            "rows_sampled": 0,
            "optimization_applied": False,
            "errors": []
        }
        
        if not PANDAS_AVAILABLE:
            error_msg = "Pandas not available for Excel processing"
            logger.error("XLSX extraction failed for %s: %s", filename, error_msg)
            metrics["errors"].append(error_msg)
            return "", metrics
        
        try:
            # Read all sheets
            sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            text_parts = []
            total_rows = 0
            
            for sheet_name, df in sheets.items():
                metrics["sheets_processed"] += 1
                original_rows = len(df)
                total_rows += original_rows
                
                # Smart processing based on sheet size
                if original_rows <= 50:
                    # Small sheet: include everything
                    sheet_text = f"=== Sheet: {sheet_name} ({original_rows} rows) ===\n"
                    sheet_text += df.to_string(max_cols=20)
                    metrics["rows_sampled"] += original_rows
                    
                elif original_rows <= 200:
                    # Medium sheet: include headers + sample
                    sheet_text = f"=== Sheet: {sheet_name} ({original_rows} rows - showing all) ===\n"
                    sheet_text += df.to_string(max_rows=200, max_cols=15)
                    metrics["rows_sampled"] += original_rows
                    
                else:
                    # Large sheet: intelligent sampling
                    metrics["optimization_applied"] = True
                    
                    # 1. Always include header info
                    sheet_text = f"=== Sheet: {sheet_name} ({original_rows} rows - sampled) ===\n"
                    sheet_text += f"Columns: {', '.join(df.columns.astype(str))}\n"
                    
                    # 2. Include first 20 rows (headers + initial data)
                    sheet_text += "\n--- First 20 rows ---\n"
                    sheet_text += df.head(20).to_string(max_cols=15) + "\n"
                    
                    # 3. Include last 10 rows (recent data)
                    sheet_text += "\n--- Last 10 rows ---\n"
                    sheet_text += df.tail(10).to_string(max_cols=15) + "\n"
                    
                    # 4. Sample middle rows if dataset is very large
                    if original_rows > 500:
                        middle_start = original_rows // 3
                        middle_sample = df.iloc[middle_start:middle_start + 15]
                        sheet_text += f"\n--- Sample from middle (rows {middle_start}-{middle_start + 15}) ---\n"
                        sheet_text += middle_sample.to_string(max_cols=15) + "\n"
                        
                    # 5. Add summary statistics for numerical columns
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    if len(numeric_cols) > 0 and len(numeric_cols) <= 10:
                        sheet_text += "\n--- Summary Statistics ---\n"
                        sheet_text += df[numeric_cols].describe().round(2).to_string() + "\n"
                    
                    # 6. Add unique value counts for categorical columns (limited)
                    categorical_cols = df.select_dtypes(include=['object']).columns
                    for col in categorical_cols[:3]:  # Limit to first 3 categorical columns
                        unique_count = df[col].nunique()
                        if unique_count <= 20:  # Only show if reasonable number of unique values
                            sheet_text += f"\n--- Unique values in '{col}' ---\n"
                            sheet_text += df[col].value_counts().head(10).to_string() + "\n"
                    
                    metrics["rows_sampled"] += 45  # 20 + 10 + 15 sampled rows
                    metrics["errors"].append(f"Sheet '{sheet_name}' sampled from {original_rows} rows to preserve key information")
                
                text_parts.append(sheet_text)
            
            # Add overall summary
            summary = f"=== Excel File Summary ===\n"
            summary += f"Total sheets: {len(sheets)}\n"
            summary += f"Total rows across all sheets: {total_rows:,}\n"
            summary += f"Rows included in extraction: {metrics['rows_sampled']:,}\n"
            if metrics["optimization_applied"]:
                summary += f"Optimization applied: Smart sampling used for large sheets\n"
            summary += f"Processing mode: {metrics['extraction_mode']}\n\n"
            
            # Combine all content
            text_content = summary + "\n\n".join(text_parts)
            
            # Final size check and truncation if needed
            if len(text_content) > 100000:  # 100KB limit
                text_content = text_content[:100000] + "\n\n... [Content truncated - use file preview for complete data]"
                metrics["errors"].append("Final content truncated to 100KB for processing efficiency")
            
            metrics["total_chars"] = len(text_content)
            logger.info("XLSX processed: %d sheets, %d total rows, %d sampled, %d chars", 
                       len(sheets), total_rows, metrics['rows_sampled'], len(text_content))
            return text_content, metrics
            
        except Exception as e:
            error_msg = f"Excel extraction failed: {str(e)}"
            logger.error("Excel extraction failed for %s: %s", filename, str(e))
            metrics["errors"].append(error_msg)
            return "", metrics
    
    @staticmethod
    def extract_json_text(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from JSON files."""
        logger.info("Processing JSON: %s", filename)
        
        metrics = {
            "extraction_mode": "Direct",
            "total_chars": 0,
            "errors": []
        }
        
        try:
            # Try UTF-8 first, fallback to latin1
            try:
                json_text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                json_text = file_bytes.decode("latin1", errors="ignore")
                metrics["errors"].append("Used latin1 encoding fallback")
            
            data = json.loads(json_text)
            
            # Pretty print with reasonable limits
            text_content = json.dumps(data, indent=2, ensure_ascii=False)
            
            # Limit very large JSON files
            if len(text_content) > 50000:  # 50KB limit
                text_content = text_content[:50000] + "\n... [Content truncated due to size]"
                metrics["errors"].append("Large JSON truncated to 50KB")
            
            metrics["total_chars"] = len(text_content)
            logger.info("JSON extracted: %d chars", len(text_content))
            return text_content, metrics
            
        except Exception as e:
            error_msg = f"JSON extraction failed: {str(e)}"
            logger.error("JSON extraction failed for %s: %s", filename, str(e))
            metrics["errors"].append(error_msg)
            return "", metrics
    
    @staticmethod
    def extract_markdown_text(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Markdown files."""
        logger.info("Processing Markdown: %s", filename)
        
        metrics = {
            "extraction_mode": "Direct",
            "total_chars": 0,
            "errors": []
        }
        
        # Try Unstructured first
        if UNSTRUCTURED_AVAILABLE:
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".md") as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                
                loader = UnstructuredMarkdownLoader(tmp_path)
                docs = loader.load()
                text_content = "\n".join([d.page_content for d in docs])
                
                # Cleanup
                os.unlink(tmp_path)
                
                metrics["total_chars"] = len(text_content)
                logger.info("Markdown extracted via Unstructured: %d chars", len(text_content))
                return text_content, metrics
                
            except Exception as e:
                logger.warning("Unstructured Markdown extraction failed for %s: %s", filename, str(e))
                metrics["errors"].append(f"Unstructured failed: {str(e)}")
        
        # Fallback to plain text
        try:
            text_content = file_bytes.decode("utf-8", errors="ignore")
            metrics["total_chars"] = len(text_content)
            logger.info("Markdown extracted as plain text: %d chars", len(text_content))
            return text_content, metrics
            
        except Exception as e:
            error_msg = f"Markdown extraction failed: {str(e)}"
            logger.error("Markdown extraction failed for %s: %s", filename, str(e))
            metrics["errors"].append(error_msg)
            return "", metrics
    
    @staticmethod
    def extract_txt_text(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from plain text files."""
        logger.info("Processing TXT: %s", filename)
        
        metrics = {
            "extraction_mode": "Direct",
            "total_chars": 0,
            "errors": []
        }
        
        try:
            # Try UTF-8 first, fallback to other encodings
            encodings = ["utf-8", "latin1", "cp1252", "ascii"]
            
            for encoding in encodings:
                try:
                    text_content = file_bytes.decode(encoding)
                    if encoding != "utf-8":
                        metrics["errors"].append(f"Used {encoding} encoding")
                    break
                except UnicodeDecodeError:
                    continue
            else:
                # If all encodings fail, use UTF-8 with errors ignored
                text_content = file_bytes.decode("utf-8", errors="ignore")
                metrics["errors"].append("Used UTF-8 with error ignoring")
            
            metrics["total_chars"] = len(text_content)
            logger.info("TXT extracted: %d chars", len(text_content))
            return text_content, metrics
            
        except Exception as e:
            error_msg = f"Text extraction failed: {str(e)}"
            logger.error("Text extraction failed for %s: %s", filename, str(e))
            metrics["errors"].append(error_msg)
            return "", metrics
    
    @staticmethod
    def extract_html_text(file_bytes: bytes, filename: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from HTML files."""
        logger.info("Processing HTML: %s", filename)
        
        metrics = {
            "extraction_mode": "Direct",
            "total_chars": 0,
            "errors": []
        }
        
        try:
            # Try BeautifulSoup if available
            try:
                from bs4 import BeautifulSoup
                
                html_content = file_bytes.decode("utf-8", errors="ignore")
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                
                text_content = soup.get_text(separator='\n', strip=True)
                metrics["total_chars"] = len(text_content)
                logger.info("HTML extracted via BeautifulSoup: %d chars", len(text_content))
                return text_content, metrics
                
            except ImportError:
                # Fallback to basic HTML tag removal
                import re
                html_content = file_bytes.decode("utf-8", errors="ignore")
                # Remove HTML tags
                text_content = re.sub(r'<[^>]+>', '', html_content)
                # Clean up whitespace
                text_content = re.sub(r'\s+', ' ', text_content).strip()
                
                metrics["total_chars"] = len(text_content)
                metrics["errors"].append("Used basic HTML tag removal (BeautifulSoup not available)")
                logger.info("HTML extracted via regex: %d chars", len(text_content))
                return text_content, metrics
                
        except Exception as e:
            error_msg = f"HTML extraction failed: {str(e)}"
            logger.error("HTML extraction failed for %s: %s", filename, str(e))
            metrics["errors"].append(error_msg)
            return "", metrics


async def process_document(
    content: str,
    content_type: str,
    filename: str,
    force_ocr: bool = False
) -> Tuple[str, int, Dict[str, Any]]:
    """
    Universal document processing pipeline.
    
    Args:
        content: Base64 encoded content
        content_type: MIME type of the document
        filename: Original filename
        force_ocr: If True, forces OCR processing for PDFs
    
    Returns:
        Tuple of (extracted_text, size_in_bytes, processing_metrics)
    """
    logger.info("Processing document: %s (type: %s)", filename, content_type)
    logger.debug("Content length: %d chars", len(content))
    
    # Initialize processing metrics
    processing_metrics = {
        "filename": filename,
        "content_type": content_type,
        "file_type": None,
        "extraction_mode": "Direct",
        "total_chars": 0,
        "processing_time_ms": 0,
        "errors": []
    }
    
    start_time = pd.Timestamp.now() if PANDAS_AVAILABLE else None
    
    try:
        # Get configuration for document limits
        await config_manager.initialize()
        config = config_manager.get_current_config()
        doc_limits = config.document_limits
        
        # Decode base64 content to check file size
        try:
            file_bytes = base64.b64decode(content)
            file_size = len(file_bytes)
            file_size_mb = file_size / (1024 * 1024)
            
            logger.debug("Decoded file size: %d bytes (%.2f MB)", file_size, file_size_mb)
            
            # Check file size limits
            if file_size_mb > doc_limits.max_file_size_mb:
                error_msg = f"File too large: {file_size_mb:.2f} MB exceeds limit of {doc_limits.max_file_size_mb} MB"
                logger.error("REJECTED: %s", error_msg)
                processing_metrics["errors"].append(error_msg)
                raise ValueError(error_msg)
            
            # Warning for large files
            if file_size_mb > doc_limits.warn_file_size_mb:
                warning_msg = f"Large file warning: {file_size_mb:.2f} MB (above {doc_limits.warn_file_size_mb} MB threshold)"
                logger.warning(warning_msg)
                processing_metrics["warnings"] = processing_metrics.get("warnings", [])
                processing_metrics["warnings"].append(warning_msg)
                
        except Exception as decode_error:
            error_msg = f"Failed to decode base64 content: {str(decode_error)}"
            logger.error(error_msg)
            processing_metrics["errors"].append(error_msg)
            return "", 0, processing_metrics
        # Detect file type
        file_type = detect_file_type(content_type, filename)
        processing_metrics["file_type"] = file_type
        logger.info("Detected file type: %s", file_type)
        
        # Route to appropriate extractor (file_bytes already decoded earlier)
        extractor = DocumentExtractor()
        
        if file_type == "pdf":
            text_content, extraction_metrics = await extractor.extract_pdf_text(
                file_bytes, filename, force_ocr
            )
        elif file_type == "docx":
            text_content, extraction_metrics = extractor.extract_docx_text(
                file_bytes, filename
            )
        elif file_type == "pptx":
            text_content, extraction_metrics = extractor.extract_pptx_text(
                file_bytes, filename
            )
        elif file_type == "csv":
            text_content, extraction_metrics = extractor.extract_csv_text(
                file_bytes, filename
            )
        elif file_type == "xlsx":
            text_content, extraction_metrics = extractor.extract_xlsx_text(
                file_bytes, filename
            )
        elif file_type == "json":
            text_content, extraction_metrics = extractor.extract_json_text(
                file_bytes, filename
            )
        elif file_type == "markdown":
            text_content, extraction_metrics = extractor.extract_markdown_text(
                file_bytes, filename
            )
        elif file_type == "txt":
            text_content, extraction_metrics = extractor.extract_txt_text(
                file_bytes, filename
            )
        elif file_type == "html":
            text_content, extraction_metrics = extractor.extract_html_text(
                file_bytes, filename
            )
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        # Merge extraction metrics
        processing_metrics.update(extraction_metrics)
        processing_metrics["total_chars"] = len(text_content)
        
        # Check extracted character count limits
        char_count = len(text_content)
        if char_count > doc_limits.max_extracted_chars:
            error_msg = f"Extracted content too large: {char_count:,} characters exceeds limit of {doc_limits.max_extracted_chars:,}"
            logger.error("REJECTED: %s", error_msg)
            processing_metrics["errors"].append(error_msg)
            raise ValueError(error_msg)
        
        # Warning for large character count
        if char_count > doc_limits.warn_extracted_chars:
            warning_msg = f"Large content warning: {char_count:,} characters (above {doc_limits.warn_extracted_chars:,} threshold)"
            logger.warning(warning_msg)
            processing_metrics["warnings"] = processing_metrics.get("warnings", [])
            processing_metrics["warnings"].append(warning_msg)
        
        # Calculate processing time
        if start_time and PANDAS_AVAILABLE:
            end_time = pd.Timestamp.now()
            processing_time_ms = int((end_time - start_time).total_seconds() * 1000)
            processing_metrics["processing_time_ms"] = processing_time_ms
        
        # Validate extracted content
        if not text_content or len(text_content.strip()) < 10:
            warning_msg = f"Minimal content extracted ({len(text_content)} chars)"
            logger.warning(warning_msg)
            processing_metrics["errors"].append(warning_msg)
        
        logger.info("Processing completed: %d chars extracted from %s file", len(text_content), file_type)
        
        return text_content, file_size, processing_metrics
        
    except Exception as e:
        error_msg = f"Document processing failed: {str(e)}"
        logger.error("Document processing failed for %s: %s", filename, str(e))
        processing_metrics["errors"].append(error_msg)
        
        import traceback
        logger.debug("Traceback: %s", traceback.format_exc())
        raise